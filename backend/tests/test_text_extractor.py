"""
Markdown 抽出の整形ルールを検証する。
画像リンクや通常リンクを平文化し、Office/PDF/Outlook も検索用テキストへ変換できることを担保する。
"""

from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch

from app.extractors.text_extractor import (
    extract_text,
    normalize_extension_filter,
    resolve_supported_extension,
    supports_content_extraction,
    supports_extension,
)


def test_supports_extension_includes_office_pdf_and_images() -> None:
    """
    対応拡張子には Office/PDF/Outlook と作図系テキスト、画像・音声のファイル名検索対象が含まれる。
    """
    assert supports_extension(Path("report.docx")) is True
    assert supports_extension(Path("sheet.xlsx")) is True
    assert supports_extension(Path("macro_sheet.xlsm")) is True
    assert supports_extension(Path("deck.pptx")) is True
    assert supports_extension(Path("memo.pdf")) is True
    assert supports_extension(Path("mail.msg")) is True
    assert supports_extension(Path("canvas.excalidraw")) is True
    assert supports_extension(Path("canvas.excalidraw.md")) is True
    assert supports_extension(Path("flow.dio")) is True
    assert supports_extension(Path("flow.dio.svg")) is True
    assert supports_extension(Path("layout.xml")) is True
    assert supports_extension(Path("photo.png")) is True
    assert supports_extension(Path("song.mp3")) is True
    assert supports_extension(Path("voice.m4a")) is True
    assert supports_extension(Path("recording.wav")) is True
    assert supports_extension(Path("archive.zip")) is False


def test_extract_text_reads_excalidraw_and_dio_as_plain_text(tmp_path: Path) -> None:
    """
    Excalidraw と draw.io 系ファイルは、JSON の値だけを検索対象へ取り込む。
    """
    excalidraw_file = tmp_path / "architecture.excalidraw"
    excalidraw_file.write_text(
        '{"type":"excalidraw","elements":[{"text":"検索導線","x":120}],"appState":{"viewBackgroundColor":"#0f172a"}}',
        encoding="utf-8",
    )

    dio_file = tmp_path / "sequence.dio"
    dio_file.write_text(
        '<mxfile><diagram>{"title":"index status filter","steps":["open","search"]}</diagram></mxfile>',
        encoding="utf-8",
    )
    dio_svg_file = tmp_path / "sequence.dio.svg"
    dio_svg_file.write_text(
        '<svg><metadata>{&quot;label&quot;:&quot;embedded drawio text&quot;,&quot;count&quot;:3}</metadata><text>ignored label</text></svg>',
        encoding="utf-8",
    )

    excalidraw_text = extract_text(excalidraw_file)
    dio_text = extract_text(dio_file)
    dio_svg_text = extract_text(dio_svg_file)

    assert "検索導線" in excalidraw_text
    assert "120" in excalidraw_text
    assert '"type"' not in excalidraw_text
    assert "index status filter" in dio_text
    assert "open" in dio_text
    assert "<mxfile>" not in dio_text
    assert "embedded drawio text" in dio_svg_text
    assert "3" in dio_svg_text
    assert "ignored label" not in dio_svg_text


def test_extract_text_reads_json_values_without_json_syntax(tmp_path: Path) -> None:
    """
    JSON ファイルはキー名や記号ではなく、値だけを検索用テキストとして取り込む。
    """
    json_file = tmp_path / "settings.json"
    json_file.write_text(
        '{"name":"alpha search","enabled":true,"items":[1,"beta"],"meta":{"owner":"team"}}',
        encoding="utf-8",
    )

    extracted = extract_text(json_file)

    assert "alpha search" in extracted
    assert "true" in extracted
    assert "1" in extracted
    assert "beta" in extracted
    assert "team" in extracted
    assert '"name"' not in extracted
    assert "{" not in extracted


def test_extract_text_reads_xml_text_without_tags(tmp_path: Path) -> None:
    """
    XML ファイルはタグ名ではなく、テキストノードの中身だけを検索用テキストとして取り込む。
    """
    xml_file = tmp_path / "layout.xml"
    xml_file.write_text(
        "<root><title>alpha layout</title><item priority='high'>beta</item><meta><count>3</count></meta></root>",
        encoding="utf-8",
    )

    extracted = extract_text(xml_file)

    assert "alpha layout" in extracted
    assert "beta" in extracted
    assert "3" in extracted
    assert "<title>" not in extracted
    assert "priority" not in extracted


def test_normalize_extension_filter_accepts_space_separated_values_without_dots() -> None:
    """
    拡張子フィルタは `md excalidraw` のようなスペース区切り・ドットなし入力も受け付ける。
    """
    normalized = normalize_extension_filter("md excalidraw dio excalidraw.md dio.svg")

    assert normalized == frozenset({".md", ".excalidraw", ".dio", ".excalidraw.md", ".dio.svg"})


def test_resolve_supported_extension_prefers_longest_match_for_compound_suffix() -> None:
    """
    複合拡張子は最長一致で判定し、`.md` と `.excalidraw.md` を区別する。
    """
    assert resolve_supported_extension(Path("note.md")) == ".md"
    assert resolve_supported_extension(Path("diagram.excalidraw.md")) == ".excalidraw.md"
    assert resolve_supported_extension(Path("flow.dio.svg")) == ".dio.svg"


def test_extract_text_reads_excalidraw_markdown_as_distinct_extension(tmp_path: Path) -> None:
    """
    `.excalidraw.md` は通常の `.md` と別拡張子として扱いつつ、本文抽出は Markdown と同様に行う。
    """
    excalidraw_markdown = tmp_path / "whiteboard.excalidraw.md"
    excalidraw_markdown.write_text("# 図\n\nリンクは [こちら](docs/spec.md)", encoding="utf-8")

    extracted = extract_text(excalidraw_markdown)

    assert "図" in extracted
    assert "こちら" in extracted


def test_extract_text_reads_docx_paragraphs_and_tables(tmp_path: Path) -> None:
    """
    DOCX は本文段落と表セルを検索用テキストとして抽出する。
    """
    from docx import Document

    document = Document()
    document.add_heading("週次レポート")
    document.add_paragraph("営業進捗を確認する")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "担当"
    table.cell(0, 1).text = "件数"
    table.cell(1, 0).text = "田中"
    table.cell(1, 1).text = "12"
    docx_file = tmp_path / "report.docx"
    document.save(docx_file)

    extracted = extract_text(docx_file)

    assert "週次レポート" in extracted
    assert "営業進捗を確認する" in extracted
    assert "担当\t件数" in extracted
    assert "田中\t12" in extracted


def test_extract_text_reads_xlsx_sheet_names_and_cell_values(tmp_path: Path) -> None:
    """
    XLSX はシート名とセル値を検索用テキストとして抽出する。
    """
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "売上"
    sheet["A1"] = "担当"
    sheet["B1"] = "金額"
    sheet["A2"] = "佐藤"
    sheet["B2"] = 125000
    xlsx_file = tmp_path / "sales.xlsx"
    workbook.save(xlsx_file)

    extracted = extract_text(xlsx_file)

    assert "[Sheet] 売上" in extracted
    assert "担当\t金額" in extracted
    assert "佐藤\t125000" in extracted


def test_extract_text_reads_xlsm_sheet_names_and_cell_values(tmp_path: Path) -> None:
    """
    XLSM も XLSX と同様にシート名とセル値を検索用テキストとして抽出する。
    """
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "在庫"
    sheet["A1"] = "品名"
    sheet["B1"] = "数量"
    sheet["A2"] = "ボルト"
    sheet["B2"] = 42
    xlsm_file = tmp_path / "inventory.xlsm"
    workbook.save(xlsm_file)

    extracted = extract_text(xlsm_file)

    assert "[Sheet] 在庫" in extracted
    assert "品名\t数量" in extracted
    assert "ボルト\t42" in extracted


def test_extract_text_reads_pptx_slide_text(tmp_path: Path) -> None:
    """
    PPTX は各スライドのテキストを検索用テキストとして抽出する。
    """
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "検索基盤"
    slide.placeholders[1].text = "並列抽出を導入する"
    pptx_file = tmp_path / "plan.pptx"
    presentation.save(pptx_file)

    extracted = extract_text(pptx_file)

    assert "[Slide 1]" in extracted
    assert "検索基盤" in extracted
    assert "並列抽出を導入する" in extracted


def test_extract_text_reads_pdf_text_via_pypdf(tmp_path: Path) -> None:
    """
    PDF は各ページの抽出結果を連結して検索用テキストへ変換する。
    """
    pdf_file = tmp_path / "memo.pdf"
    pdf_file.write_bytes(b"%PDF-1.4")

    class FakePage:
        def __init__(self, text: str) -> None:
            self._text = text

        def extract_text(self, *args, **kwargs) -> str:
            return self._text

    class FakePdfReader:
        def __init__(self, *args, **kwargs) -> None:
            self.pages = [FakePage("1ページ目"), FakePage("2ページ目")]

    with patch.dict("sys.modules", {"pypdf": SimpleNamespace(PdfReader=FakePdfReader)}):
        extracted = extract_text(pdf_file)

    assert "[Page 1]" in extracted
    assert "1ページ目" in extracted
    assert "2ページ目" in extracted


def test_extract_text_skips_encrypted_pdf_without_error(tmp_path: Path) -> None:
    """
    暗号化された PDF は本文抽出を行わず、空文字を返して静かにスキップする。
    """
    pdf_file = tmp_path / "secret.pdf"
    pdf_file.write_bytes(b"%PDF-1.4")

    class FakePdfReader:
        def __init__(self, *args, **kwargs) -> None:
            self.is_encrypted = True
            self.pages = []

    with patch.dict("sys.modules", {"pypdf": SimpleNamespace(PdfReader=FakePdfReader)}):
        extracted = extract_text(pdf_file)

    assert extracted == ""


def test_extract_text_keeps_extractable_pdf_pages_when_one_page_breaks(tmp_path: Path) -> None:
    """
    PDF の一部ページで抽出エラーが起きても、読めたページの本文は取り込む。
    """
    pdf_file = tmp_path / "partial.pdf"
    pdf_file.write_bytes(b"%PDF-1.4")

    class FakePage:
        def __init__(self, *, layout_text: str | None = None, error: Exception | None = None) -> None:
            self._layout_text = layout_text
            self._error = error

        def extract_text(self, *args, **kwargs) -> str:
            if self._error is not None:
                raise self._error
            return self._layout_text or ""

    class FakePdfReader:
        def __init__(self, *args, **kwargs) -> None:
            self.is_encrypted = False
            self.pages = [
                FakePage(layout_text="読める1ページ目"),
                FakePage(error=ValueError("unknown encoding: /90ms-RKSJ-H")),
                FakePage(error=ZeroDivisionError("float floor division by zero")),
                FakePage(layout_text="読める4ページ目"),
            ]

    with patch.dict("sys.modules", {"pypdf": SimpleNamespace(PdfReader=FakePdfReader)}):
        extracted = extract_text(pdf_file)

    assert "読める1ページ目" in extracted
    assert "読める4ページ目" in extracted
    assert "90ms-RKSJ-H" not in extracted


def test_extract_text_reads_msg_header_and_body(tmp_path: Path) -> None:
    """
    Outlook の MSG は件名や差出人と本文を検索用テキストへ変換する。
    """
    msg_file = tmp_path / "mail.msg"
    msg_file.write_bytes(b"msg")

    class FakeMessage:
        subject = "見積確認"
        sender = "sales@example.com"
        to = "team@example.com"
        cc = None
        date = "2026-04-13 10:30:00"
        body = "案件の見積を確認してください。"

        def close(self) -> None:
            return None

    def fake_open_msg(path: str, **kwargs) -> FakeMessage:
        assert path.endswith("mail.msg")
        return FakeMessage()

    with patch.dict("sys.modules", {"extract_msg": SimpleNamespace(openMsg=fake_open_msg)}):
        extracted = extract_text(msg_file)

    assert "Subject: 見積確認" in extracted
    assert "From: sales@example.com" in extracted
    assert "To: team@example.com" in extracted
    assert "案件の見積を確認してください。" in extracted


def test_extract_text_returns_empty_string_when_msg_decode_fails(tmp_path: Path) -> None:
    """
    MSG の文字コード解釈に失敗した場合は、空文字へフォールバックして継続する。
    """
    msg_file = tmp_path / "broken.msg"
    msg_file.write_bytes(b"msg")

    def fake_open_msg(path: str, **kwargs):
        assert path.endswith("broken.msg")
        raise UnicodeDecodeError("iso2022_jp", b"header", 1, 2, "illegal multibyte sequence")

    with patch.dict("sys.modules", {"extract_msg": SimpleNamespace(openMsg=fake_open_msg)}):
        extracted = extract_text(msg_file)

    assert extracted == ""


def test_extract_text_reads_msg_body_when_optional_headers_fail_to_decode(tmp_path: Path) -> None:
    """
    MSG の任意ヘッダ取得で文字コード例外が出ても、件名と本文を優先して抽出し続ける。
    """
    msg_file = tmp_path / "partial.msg"
    msg_file.write_bytes(b"msg")

    class FakeMessage:
        subject = "障害連絡"
        body = "本文だけでも検索できれば十分です。"

        @property
        def sender(self) -> str:
            raise UnicodeDecodeError("cp932", b"from", 0, 1, "invalid start byte")

        @property
        def to(self) -> str:
            raise UnicodeDecodeError("cp932", b"to", 0, 1, "invalid start byte")

        cc = None
        date = None

        def close(self) -> None:
            return None

    def fake_open_msg(path: str, **kwargs) -> FakeMessage:
        assert path.endswith("partial.msg")
        return FakeMessage()

    with patch.dict("sys.modules", {"extract_msg": SimpleNamespace(openMsg=fake_open_msg)}):
        extracted = extract_text(msg_file)

    assert "Subject: 障害連絡" in extracted
    assert "From: -" in extracted
    assert "To: -" in extracted
    assert "Cc: -" in extracted
    assert "Date: -" in extracted
    assert "本文だけでも検索できれば十分です。" in extracted


def test_extract_text_returns_empty_string_when_pptx_package_is_broken(tmp_path: Path) -> None:
    """
    壊れた PPTX は空文字へフォールバックし、全体のインデックス処理を止めない。
    """
    pptx_file = tmp_path / "broken.pptx"
    pptx_file.write_bytes(b"pptx")

    def fake_presentation(path: str):
        assert path.endswith("broken.pptx")
        raise ValueError('Package not found at "/tmp/broken.pptx"')

    with patch.dict("sys.modules", {"pptx": SimpleNamespace(Presentation=fake_presentation)}):
        extracted = extract_text(pptx_file)

    assert extracted == ""


def test_extract_text_flattens_markdown_image_path_with_parentheses(tmp_path: Path) -> None:
    """
    Markdown 画像リンクの宛先に ASCII の丸括弧が含まれても、全文を保持して抽出する。
    """
    markdown_file = tmp_path / "sample.md"
    markdown_file.write_text(
        r"パーツの新規作成(流用) ![](C:\M222345\000_work\tehai_s\パーツの新規作成(flow).png)",
        encoding="utf-8",
    )

    extracted = extract_text(markdown_file)

    assert r"C:\M222345\000_work\tehai_s\パーツの新規作成(flow).png" in extracted
    assert "![](" not in extracted


def test_extract_text_keeps_plain_text_files_unchanged(tmp_path: Path) -> None:
    """
    Markdown 以外のテキストはそのまま返す。
    """
    text_file = tmp_path / "sample.txt"
    text_file.write_text("alpha ![](beta(gamma).png)", encoding="utf-8")

    assert extract_text(text_file) == "alpha ![](beta(gamma).png)"


def test_custom_content_extensions_are_treated_as_plain_text(tmp_path: Path) -> None:
    """
    利用者追加の本文抽出拡張子はプレーンテキストとして扱える。
    """
    custom_file = tmp_path / "solver.dat"
    custom_file.write_text("NODE 1001", encoding="utf-8")

    assert resolve_supported_extension(custom_file, extra_content_extensions=(".dat",)) == ".dat"
    assert supports_content_extraction(custom_file, extra_content_extensions=(".dat",)) is True
    assert "NODE 1001" in extract_text(custom_file, extra_content_extensions=(".dat",))


def test_normalize_extension_filter_accepts_custom_extension_lists() -> None:
    """
    カスタム拡張子を追加した場合も、拡張子フィルタで正規化して扱える。
    """
    normalized = normalize_extension_filter(
        "md py cae",
        extra_content_extensions=(".py",),
        extra_filename_extensions=(".cae",),
    )

    assert normalized == frozenset({".md", ".py", ".cae"})
