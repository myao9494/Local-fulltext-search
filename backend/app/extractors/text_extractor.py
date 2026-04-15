"""
検索用テキスト抽出ユーティリティ。
本文を抽出できる拡張子と、ファイル名だけを検索対象にする拡張子を分けて扱う。
"""

from __future__ import annotations

import html
import json
from datetime import date, datetime, time
from pathlib import Path
import re
from xml.etree import ElementTree


CONTENT_EXTENSIONS: frozenset[str] = frozenset(
    {
        ".md",
        ".excalidraw.md",
        ".dio.svg",
        ".json",
        ".xml",
        ".txt",
        ".excalidraw",
        ".dio",
        ".pdf",
        ".docx",
        ".xlsx",
        ".xlsm",
        ".pptx",
        ".msg",
    }
)

FILENAME_ONLY_EXTENSIONS: frozenset[str] = frozenset(
    {
        ".png",
        ".jpg",
        ".jpeg",
        ".gif",
        ".webp",
        ".heic",
        ".svg",
        ".bmp",
        ".tif",
        ".tiff",
        ".mp3",
        ".m4a",
        ".aac",
        ".wav",
        ".flac",
        ".aif",
        ".aiff",
        ".alac",
        ".m4p",
    }
)

SUPPORTED_EXTENSIONS: frozenset[str] = CONTENT_EXTENSIONS | FILENAME_ONLY_EXTENSIONS
SORTED_SUPPORTED_EXTENSIONS: tuple[str, ...] = tuple(sorted(SUPPORTED_EXTENSIONS, key=len, reverse=True))


def get_content_extensions(*, extra_content_extensions: tuple[str, ...] | list[str] | set[str] | frozenset[str] = ()) -> frozenset[str]:
    """
    標準の本文抽出拡張子に、利用者追加分を合成して返す。
    """
    return frozenset({*CONTENT_EXTENSIONS, *extra_content_extensions})


def get_filename_only_extensions(
    *, extra_filename_extensions: tuple[str, ...] | list[str] | set[str] | frozenset[str] = ()
) -> frozenset[str]:
    """
    標準のファイル名検索専用拡張子に、利用者追加分を合成して返す。
    """
    return frozenset({*FILENAME_ONLY_EXTENSIONS, *extra_filename_extensions})


def get_supported_extensions(
    *,
    extra_content_extensions: tuple[str, ...] | list[str] | set[str] | frozenset[str] = (),
    extra_filename_extensions: tuple[str, ...] | list[str] | set[str] | frozenset[str] = (),
) -> frozenset[str]:
    """
    標準拡張子と利用者追加拡張子を合わせた全対応拡張子を返す。
    """
    return get_content_extensions(extra_content_extensions=extra_content_extensions) | get_filename_only_extensions(
        extra_filename_extensions=extra_filename_extensions
    )


def normalize_extension_token(value: str) -> str:
    """
    拡張子入力を `.md` 形式へそろえる。ドットなし入力も受け付ける。
    """
    token = value.strip().lower()
    if not token:
        return ""
    return token if token.startswith(".") else f".{token}"


def supports_extension(path: Path) -> bool:
    """
    検索対象として扱う拡張子かどうかを返す。
    """
    return resolve_supported_extension(path) is not None


def supports_content_extraction(
    path: Path,
    *,
    extra_content_extensions: tuple[str, ...] | list[str] | set[str] | frozenset[str] = (),
    extra_filename_extensions: tuple[str, ...] | list[str] | set[str] | frozenset[str] = (),
) -> bool:
    """
    本文抽出に対応する拡張子かどうかを返す。
    """
    resolved = resolve_supported_extension(
        path,
        extra_content_extensions=extra_content_extensions,
        extra_filename_extensions=extra_filename_extensions,
    )
    return resolved in get_content_extensions(extra_content_extensions=extra_content_extensions)


def resolve_supported_extension(
    path: Path,
    *,
    extra_content_extensions: tuple[str, ...] | list[str] | set[str] | frozenset[str] = (),
    extra_filename_extensions: tuple[str, ...] | list[str] | set[str] | frozenset[str] = (),
) -> str | None:
    """
    対応拡張子のうち、ファイル名末尾に最長一致するものを返す。
    `excalidraw.md` のような複合拡張子を `.md` と区別する。
    """
    lower_name = path.name.lower()
    sorted_supported_extensions = tuple(
        sorted(
            get_supported_extensions(
                extra_content_extensions=extra_content_extensions,
                extra_filename_extensions=extra_filename_extensions,
            ),
            key=len,
            reverse=True,
        )
    )
    for extension in sorted_supported_extensions:
        if lower_name.endswith(extension):
            return extension
    return None


def normalize_extension_filter(
    value: str | None,
    *,
    extra_content_extensions: tuple[str, ...] | list[str] | set[str] | frozenset[str] = (),
    extra_filename_extensions: tuple[str, ...] | list[str] | set[str] | frozenset[str] = (),
) -> frozenset[str]:
    """
    利用者が指定した拡張子一覧を正規化し、未指定時は全対応拡張子を返す。
    """
    supported_extensions = get_supported_extensions(
        extra_content_extensions=extra_content_extensions,
        extra_filename_extensions=extra_filename_extensions,
    )
    if value is None or not value.strip():
        return supported_extensions

    normalized = {
        normalize_extension_token(item)
        for item in re.split(r"[\s,]+", value.strip())
        if item.strip()
    }
    filtered = {item for item in normalized if item in supported_extensions}
    return frozenset(filtered) if filtered else supported_extensions


def extract_text(
    path: Path,
    *,
    extra_content_extensions: tuple[str, ...] | list[str] | set[str] | frozenset[str] = (),
    extra_filename_extensions: tuple[str, ...] | list[str] | set[str] | frozenset[str] = (),
) -> str:
    """
    ファイル種別に応じて検索用テキストを抽出する。
    Markdown は `![](...)` / `[](...)` を平文化し、Office/PDF/Outlook は平文へ変換する。
    """
    suffix = path.suffix.lower()
    resolved_extension = resolve_supported_extension(
        path,
        extra_content_extensions=extra_content_extensions,
        extra_filename_extensions=extra_filename_extensions,
    )
    if resolved_extension is None:
        raise ValueError(f"Unsupported extension: {suffix}")

    if resolved_extension in {".md", ".excalidraw.md"}:
        content = path.read_text(encoding="utf-8", errors="ignore")
        return _flatten_markdown_inline_links(content)
    if resolved_extension == ".txt":
        return path.read_text(encoding="utf-8", errors="ignore")
    if resolved_extension in {".json", ".excalidraw"}:
        return _extract_json_values_text(path.read_text(encoding="utf-8", errors="ignore"))
    if resolved_extension in {".dio", ".dio.svg"}:
        return _extract_embedded_json_values_text(path.read_text(encoding="utf-8", errors="ignore"))
    if resolved_extension == ".xml":
        return _extract_xml_text(path.read_text(encoding="utf-8", errors="ignore"))
    if resolved_extension in extra_content_extensions:
        return path.read_text(encoding="utf-8", errors="ignore")
    if resolved_extension == ".pdf":
        return _extract_pdf_text(path)
    if resolved_extension == ".docx":
        return _extract_docx_text(path)
    if resolved_extension in {".xlsx", ".xlsm"}:
        return _extract_xlsx_text(path)
    if resolved_extension == ".pptx":
        return _extract_pptx_text(path)
    if resolved_extension == ".msg":
        return _extract_msg_text(path)

    raise ValueError(f"Unsupported extension: {resolved_extension}")


def _extract_json_values_text(content: str) -> str:
    """
    JSON 文字列を解析し、キー名や記号を除いた値だけを改行区切りで返す。
    """
    try:
        parsed = json.loads(content)
    except json.JSONDecodeError:
        return ""
    return "\n".join(_iter_json_scalar_values(parsed))


def _extract_embedded_json_values_text(content: str) -> str:
    """
    XML や SVG に埋め込まれた JSON 断片を走査し、見つかった値だけを返す。
    """
    fragments = _extract_json_values_text(content)
    if fragments:
        return fragments

    values: list[str] = []
    for fragment in _iter_json_fragments(html.unescape(content)):
        values.extend(_iter_json_scalar_values(fragment))
    return "\n".join(values)


def _extract_xml_text(content: str) -> str:
    """
    XML を解析し、タグや属性を除いたテキストノードだけを改行区切りで返す。
    """
    try:
        root = ElementTree.fromstring(content)
    except ElementTree.ParseError:
        return ""

    values = [text.strip() for text in root.itertext() if text.strip()]
    return "\n".join(values)


def _iter_json_scalar_values(value: object) -> list[str]:
    """
    JSON オブジェクトを再帰的にたどり、検索対象にしたい末端値だけを抽出する。
    """
    if value is None:
        return []
    if isinstance(value, dict):
        values: list[str] = []
        for nested in value.values():
            values.extend(_iter_json_scalar_values(nested))
        return values
    if isinstance(value, list):
        values: list[str] = []
        for nested in value:
            values.extend(_iter_json_scalar_values(nested))
        return values
    if isinstance(value, bool):
        return ["true" if value else "false"]
    if isinstance(value, (int, float)):
        return [str(value)]
    if isinstance(value, str):
        stripped = value.strip()
        return [stripped] if stripped else []
    return [str(value)]


def _iter_json_fragments(content: str) -> list[object]:
    """
    テキスト中の JSON オブジェクト・配列断片を抽出して返す。
    """
    parsed_fragments: list[object] = []
    length = len(content)
    index = 0
    while index < length:
        char = content[index]
        if char not in "{[":
            index += 1
            continue

        fragment, next_index = _read_json_fragment(content, index)
        if fragment is None:
            index += 1
            continue

        try:
            parsed_fragments.append(json.loads(fragment))
        except json.JSONDecodeError:
            index += 1
            continue

        index = next_index

    return parsed_fragments


def _read_json_fragment(content: str, start_index: int) -> tuple[str | None, int]:
    """
    開始位置から JSON らしい波括弧/角括弧ブロックを切り出す。
    """
    opening = content[start_index]
    closing = "}" if opening == "{" else "]"
    stack = [closing]
    index = start_index + 1
    in_string = False
    is_escaped = False

    while index < len(content):
        char = content[index]
        if in_string:
            if is_escaped:
                is_escaped = False
            elif char == "\\":
                is_escaped = True
            elif char == '"':
                in_string = False
            index += 1
            continue

        if char == '"':
            in_string = True
            index += 1
            continue
        if char == "{":
            stack.append("}")
            index += 1
            continue
        if char == "[":
            stack.append("]")
            index += 1
            continue
        if stack and char == stack[-1]:
            stack.pop()
            index += 1
            if not stack:
                return content[start_index:index], index
            continue
        index += 1

    return None, start_index


def _extract_pdf_text(path: Path) -> str:
    """
    PDF の各ページからテキストを取り出し、ページ単位で連結する。
    暗号化 PDF は無理に復号せずスキップし、ページ単位の抽出失敗は読めたページだけを残す。
    """
    try:
        from pypdf import PdfReader
    except ImportError as error:
        raise RuntimeError("PDF 抽出には pypdf が必要です。") from error

    reader = PdfReader(str(path), strict=False)
    if bool(getattr(reader, "is_encrypted", False)):
        return ""

    parts: list[str] = []
    for page_index, page in enumerate(reader.pages, start=1):
        text = _extract_pdf_page_text(page)
        stripped = text.strip()
        if stripped:
            parts.append(f"[Page {page_index}]\n{stripped}")
    return "\n\n".join(parts)


def _extract_pdf_page_text(page) -> str:
    """
    PDF 1ページ分の本文を抽出する。
    layout 抽出に失敗しても通常抽出へフォールバックし、既知の不正PDF系エラーは空文字扱いにする。
    """
    for kwargs in ({"extraction_mode": "layout"}, {}):
        try:
            return page.extract_text(**kwargs) or ""
        except Exception as error:
            if _is_recoverable_pdf_error(error):
                continue
            raise
    return ""


def _is_recoverable_pdf_error(error: Exception) -> bool:
    """
    暗号化や壊れた埋め込みフォントなど、既知の PDF 抽出失敗はスキップ可能と判定する。
    """
    if isinstance(error, ZeroDivisionError):
        return True
    message = str(error)
    recoverable_markers = (
        "unknown encoding:",
        "File has not been decrypted",
        "float floor division by zero",
    )
    return any(marker in message for marker in recoverable_markers)


def _extract_docx_text(path: Path) -> str:
    """
    DOCX の段落と表を文書順で平文化する。
    """
    try:
        from docx import Document
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError as error:
        raise RuntimeError("DOCX 抽出には python-docx が必要です。") from error

    document = Document(str(path))
    parts: list[str] = []
    for block in document.iter_inner_content():
        if isinstance(block, Paragraph):
            text = block.text.strip()
            if text:
                parts.append(text)
            continue
        if isinstance(block, Table):
            table_text = _extract_docx_table_text(block)
            if table_text:
                parts.append(table_text)
    return "\n\n".join(parts)


def _extract_docx_table_text(table) -> str:
    """
    DOCX 表を行ごとのタブ区切りテキストへ変換する。
    """
    rows: list[str] = []
    for row in table.rows:
        values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
        if values:
            rows.append("\t".join(values))
    return "\n".join(rows)


def _extract_xlsx_text(path: Path) -> str:
    """
    XLSX / XLSM のシート名とセル値を読み取り専用で抽出する。
    """
    try:
        from openpyxl import load_workbook
    except ImportError as error:
        raise RuntimeError("XLSX / XLSM 抽出には openpyxl が必要です。") from error

    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    try:
        parts: list[str] = []
        for sheet in workbook.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                values = [_stringify_cell_value(value) for value in row]
                present_values = [value for value in values if value]
                if present_values:
                    rows.append("\t".join(present_values))
            if rows:
                parts.append(f"[Sheet] {sheet.title}\n" + "\n".join(rows))
            else:
                parts.append(f"[Sheet] {sheet.title}")
        return "\n\n".join(parts)
    finally:
        workbook.close()


def _stringify_cell_value(value: object) -> str:
    """
    Excel のセル値を検索向け文字列へ正規化する。
    """
    if value is None:
        return ""
    if isinstance(value, datetime):
        return value.isoformat(sep=" ")
    if isinstance(value, (date, time)):
        return value.isoformat()
    return str(value).strip()


def _extract_pptx_text(path: Path) -> str:
    """
    PPTX の各スライドからテキストフレームと表を抽出する。
    """
    try:
        from pptx import Presentation
    except ImportError as error:
        raise RuntimeError("PPTX 抽出には python-pptx が必要です。") from error

    try:
        presentation = Presentation(str(path))
    except Exception as error:
        if _is_recoverable_pptx_error(error):
            return ""
        raise
    slides: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                table_text = _extract_pptx_table_text(shape.table)
                if table_text:
                    parts.append(table_text)
        if parts:
            slides.append(f"[Slide {slide_index}]\n" + "\n\n".join(parts))
    return "\n\n".join(slides)


def _extract_pptx_table_text(table) -> str:
    """
    PPTX 表を行ごとのタブ区切りテキストへ変換する。
    """
    rows: list[str] = []
    for row in table.rows:
        values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
        if values:
            rows.append("\t".join(values))
    return "\n".join(rows)


def _extract_msg_text(path: Path) -> str:
    """
    Outlook の MSG から件名・差出人・宛先・本文を抽出する。
    """
    try:
        import extract_msg
    except ImportError as error:
        raise RuntimeError("MSG 抽出には extract-msg が必要です。") from error

    try:
        message = extract_msg.openMsg(str(path), delayAttachments=True)
    except UnicodeDecodeError:
        return ""
    try:
        headers = [
            ("Subject", getattr(message, "subject", None)),
            ("From", getattr(message, "sender", None)),
            ("To", getattr(message, "to", None)),
            ("Cc", getattr(message, "cc", None)),
            ("Date", getattr(message, "date", None)),
        ]
        parts = [f"{label}: {value}".strip() for label, value in headers if value]
        body = str(getattr(message, "body", "") or "").strip()
        if body:
            parts.append(body)
        return "\n".join(parts)
    finally:
        close = getattr(message, "close", None)
        if callable(close):
            close()


def _is_recoverable_pptx_error(error: Exception) -> bool:
    """
    壊れた zip/package 由来で開けない PPTX は、本文抽出不能として空文字へフォールバックする。
    """
    return "Package not found at" in str(error)


def _flatten_markdown_inline_links(content: str) -> str:
    """
    Markdown のインライン画像/リンク記法を平文へ変換する。
    宛先パス中の括弧を深さ付きで読むため、`(flow)` のような文字を含んでも途中で切らない。
    """
    fragments: list[str] = []
    index = 0
    length = len(content)

    while index < length:
        parsed = _parse_markdown_inline_link(content, index)
        if parsed is None:
            fragments.append(content[index])
            index += 1
            continue

        label, destination, next_index = parsed
        replacement = " ".join(part for part in (label.strip(), destination.strip()) if part.strip())
        fragments.append(replacement)
        index = next_index

    return "".join(fragments)


def _parse_markdown_inline_link(content: str, start_index: int) -> tuple[str, str, int] | None:
    """
    指定位置が Markdown の `![](...)` または `[](...)` なら、その内容を返す。
    """
    index = start_index
    if content[index] == "!":
        index += 1
        if index >= len(content) or content[index] != "[":
            return None
    elif content[index] != "[":
        return None

    label, next_index = _read_bracket_content(content, index, "[", "]")
    if label is None or next_index >= len(content) or content[next_index] != "(":
        return None

    destination, end_index = _read_parenthesized_destination(content, next_index)
    if destination is None:
        return None

    return label, destination, end_index


def _read_bracket_content(
    content: str,
    start_index: int,
    opening: str,
    closing: str,
) -> tuple[str | None, int]:
    r"""
    単純な括弧ブロックを読み取る。
    `\]` のような Markdown エスケープだけを特別扱いし、Windows パスの `\` は保持する。
    """
    if content[start_index] != opening:
        return None, start_index

    index = start_index + 1
    parts: list[str] = []
    while index < len(content):
        char = content[index]
        if char == "\\" and index + 1 < len(content):
            escaped_char = content[index + 1]
            if escaped_char in {opening, closing, "\\"}:
                parts.append(escaped_char)
                index += 2
                continue
        if char == closing:
            return "".join(parts), index + 1
        parts.append(char)
        index += 1
    return None, start_index


def _read_parenthesized_destination(content: str, start_index: int) -> tuple[str | None, int]:
    r"""
    宛先の丸括弧を深さ付きで読み取る。
    `C:/tmp/a(test).png` のような ASCII 括弧を含むパスも正しく最後まで取得する。
    `\)` などの Markdown エスケープのみ展開し、Windows パスの `\` は維持する。
    """
    if content[start_index] != "(":
        return None, start_index

    index = start_index + 1
    depth = 1
    parts: list[str] = []
    while index < len(content):
        char = content[index]
        if char == "\\" and index + 1 < len(content):
            escaped_char = content[index + 1]
            if escaped_char in {"(", ")", "\\"}:
                parts.append(escaped_char)
                index += 2
                continue
        if char == "(":
            depth += 1
            parts.append(char)
            index += 1
            continue
        if char == ")":
            depth -= 1
            if depth == 0:
                return "".join(parts), index + 1
            parts.append(char)
            index += 1
            continue
        parts.append(char)
        index += 1

    return None, start_index
